import functools
import os
import sys
import shutil
import string
import uuid
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from tabulate import tabulate
from bs4 import BeautifulSoup
from langchain.docstore.document import Document
from langchain.document_loaders import UnstructuredFileLoader

import open_ai_util
import paragraph_parser
import supabase_docstore as custom_supabase
from chroma_db_util import ChromaDB

from dotenv import load_dotenv
load_dotenv()


image_count = 0
image_folder = 'ppt_images'
documents_table = os.environ.get('SUPABASE_DOCUMENTS_TABLE', '')
embeddings_table = os.environ.get('SUPABASE_EMBEDDINGS_TABLE', '')
store = custom_supabase.SupabaseDocstore(documents_table)


def pptx_table_to_html(table):
    """
    Accepts a Table object from the pptx and returns a Html table generated from it
    """
    rows = []
    for row in table.rows:
        cells = [cell.text for cell in row.cells]
        rows.append(cells)

    html_table = tabulate(rows, headers='firstrow', tablefmt='html')
    return html_table


def validate_intersection(cord1, cord2):
    """
    Takes cordinates for two lines and returns true if the more than 80% of the smaller line 
    overlaps on the longer line. 
    """
    l1, r1 = cord1
    l2, r2 = cord2
    valid_intersection = False
    intersection_length = max(0, min(r1, r2)-max(l1, l2))
    if intersection_length:
        # We need to validate percentage here?
        smaller_side = min(r1-l1, r2-l2)+1
        percentage_overlap = intersection_length * 100 / smaller_side
        valid_intersection = percentage_overlap > 80
    return valid_intersection


def compare(item1, item2):
    """
    Sorts the items in a ppt based on creation order, bottom right corner, and basic layout parsing
    """
    index1, slide1 = item1
    index2, slide2 = item2
    # Pick the element who's cords end earlier
    left1, top1 = slide1.left+slide1.width, slide1.top+slide2.height
    left2, top2 = slide2.left+slide2.width, slide2.top+slide2.height
    if left1 <= left2 and top1 <= top2:
        return -1
    elif left1 >= left2 and top1 >= top2:
        return 1
    else:
        # HANDLE VERTICAL CASES
        vertical_intersection = validate_intersection(
            (slide1.left, left1), (slide2.left, left2))
        horizontal_intersection = validate_intersection(
            (slide1.top, top1), (slide2.top, top2))
        if vertical_intersection:
            return top1-top2
        if horizontal_intersection:
            return left1-left2

        if left1 <= left2:
            # FIRST IS ON LEFT OF SECOND and top1 > top2
            # LETS RETURN FIRST AS OF NOW
            return -1
        else:
            return 1

        # FOLLOWING CODE WILL NEVER BE EXECUTED

        # IF NOTHING RESOLVES IT, RETURN ON THE BASIS ON INSERTION ORDER ITSELF
        return index1 - index2

        # HANDLE HORIZONTAL CASES


def preprocess_slide_layout(slide):
    """
    Perfroms the necessary tranformation on the processing order of shapes
    """
    sorted_shapes = [(index, shape)
                     for index, shape in enumerate(slide.shapes)]
    sorted_shapes.sort(key=functools.cmp_to_key(compare))
    return sorted_shapes


def process_shapes(slide, page_number, skip_text=None):
    """
    Extracts data from the slide/Group Shapes and returns a List of Document(s)
    """
    global image_count
    if not os.path.exists(image_folder):
        os.makedirs(image_folder)
    parsed_docs = []
    run_text = ""

    # THIS HELPS TO MAINTAIN ORDER OF TEXT
    def append_run_text(run_text):
        if run_text:
            parsed_docs.append(Document(page_content=run_text, metadata={
                'category': 'text', 'page_number': page_number}))
        return ""

    sorted_shapes = preprocess_slide_layout(slide)
    for _, shape in sorted_shapes:
        metadata = {
            'category': shape.shape_type,
            'page_number': page_number
        }
        page_content = ''
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            run_text = append_run_text(run_text)
            # TODO:: Handle cases when the image has charts,tables etc
            page_content = ''  # This could be updated if needed
            metadata['image_path'] = f'{image_folder}/img_{image_count}.jpeg'
            # metadata['cordinates'] # TODO:: Add support for narrative text via coridinates if possible
            with open(metadata['image_path'], 'wb') as f:
                f.write(shape.image.blob)
                # TODO:: VERIFY IF PARSING IMAGE VIA UNSTRUCTURED HERE IS MORE REASONABLE
                image_count += 1
            parsed_docs.append(Document(page_content='', metadata=metadata))
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # this will add the groups as parent doc as well
            run_text = append_run_text(run_text)
            group_list = process_shapes(shape, page_number)
            print('GROUP--> ', group_list)
            parsed_docs.extend(group_list)
        if shape.has_table:
            run_text = append_run_text(run_text)
            # Concvert Table object to html table
            metadata['text_as_html'] = pptx_table_to_html(shape.table)
            parsed_docs.append(Document(page_content='', metadata=metadata))
            # Now this table will function same as unstructured's table
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                # is a part of the text run.text
                text = run.text
                if skip_text:
                    text = text.replace(skip_text, "")
                if text:
                    run_text += f' {text.strip()}'
    append_run_text(run_text)
    return parsed_docs


def process_title_for_slide(slide, last_title):
    """
    Checks if the curr title and the last title are part of a single message.
    """
    curr_title = slide.shapes.title
    update_title = True
    skip_text = None
    if curr_title:
        skip_text = curr_title.text

        if last_title:
            # CASE 1: last_title covers over 70 percent of the curr title as complete prefix
            if curr_title.text.startswith(last_title.rstrip(string.punctuation + string.whitespace)) and 0.7 * len(curr_title.text) < len(last_title):
                # matches over 70 % from the last title as complete prefix
                update_title = False

            # CASE 2: Last title covers over 85% percent of the curr title but not as complete prefix
            if curr_title.text.startswith(last_title.rstrip(string.punctuation + string.whitespace)[0:int(0.85*len(last_title))]):
                update_title = False

        # append this to list as well
        if update_title:
            last_title = curr_title.text
    return (skip_text, last_title)


def process_pptx_files(filepath):
    """
    Process a pptx file and returns a List of Documents(s)
    """
    f = open(f'{filepath}', "rb")
    prs = Presentation(f)
    docs = []
    page_number = 1
    last_title = None
    for slide in prs.slides:
        skip_text, new_title = process_title_for_slide(slide, last_title)
        if new_title and new_title != last_title:
            docs.append(Document(page_content=new_title, metadata={
                        'page_number': page_number,
                        'category': 'Title'
                        }))
            last_title = new_title
        docs.extend(process_shapes(slide, page_number, skip_text))
        page_number += 1
    return docs


def process_ppt_via_unstructured(file_path):
    """
    Process ppt files via UnstructuredFileLoader
    # CURRENTLY NOT IN USE
    """
    docs = UnstructuredFileLoader(
        file_path=file_path,
        strategy='hi_res',
        mode='elements',
        include_slide_notes=True,  # WILL INCLUDE READER NOTES AS WELL
        include_page_breaks=True,
        skip_infer_table_types=[],
        pdf_infer_table_structure=True,
        pdf_extract_images=True,
        chunking_strategy='by-title'
    ).load()
    parent_doc = None
    last_narrative_text = None
    last_header = None
    child_doc_list = []
    parent_doc_list = []
    doc_uid = str(uuid.uuid4())
    for document in docs:
        page_content = document.page_content
        metadata = document.metadata
        category = metadata['category']
        table_detals = metadata.get('text_as_html')
        metadata['id'] = doc_uid
        if category == 'Title':
            last_header = page_content
        if category in ['Text', 'NarrativeText']:
            last_narrative_text = page_content
        if category == 'PageBreak':
            # create indexing here
            parent_doc = None
            continue
        if table_detals:
            table_uuid = str(uuid.uuid4())
            metadata['id'] = table_uuid
            parent_doc_list.append((
                metadata['id'],
                Document(
                    page_content=table_detals, metadata=metadata)
            ))
            child_doc_list.extend(open_ai_util.get_table_description(
                table=table_detals,
                header=last_header,
                narrative_text=last_narrative_text,
                metadata=metadata
            ))
        parent_doc += page_content


def preprocess_image_on_size_and_words(image_path):
    size_in_kb = os.stat(image_path).st_size/1024
    image_data = None
    word_count = None
    if size_in_kb < 50:
        image_data = UnstructuredFileLoader(
            file_path=image_path,
            strategy='hi_res',
            # mode='elements',
            include_slide_notes=True,  # WILL INCLUDE READER NOTES AS WELL
            include_page_breaks=True,
            skip_infer_table_types=[],
            pdf_infer_table_structure=True,
            pdf_extract_images=True,
            chunking_strategy='by-title'
        ).load()[0].page_content
        word_count = paragraph_parser.get_word_count(image_data)
        if word_count < 5 or word_count > 200:
            return image_data


def process_image_via_llm(image_path, metadata):
    """
    Extracted data from image and generates a summary for it
    """
    page_content = open_ai_util.get_image_description(image_path)
    parent_doc_text = None
    child_doc_list = []
    # CHECK FOR A VALID RESPONSE
    if page_content not in ['UNEXTRACTABLE_DATA', 'UNPROCESSABLE_ENTITY']:
        # image_id = str(uuid.uuid4())
        # # use the same metadata as of the parent image
        # metadata['id'] = image_id
        parent_doc_text = page_content
        child_doc_list.extend(
            paragraph_parser.parse_paragraph(page_content, metadata))
        # GENERATING QUESTIONS BASED ON
        child_doc_list.extend(
            open_ai_util.get_paragraph_description(paragraph=page_content, metadata=metadata))

    return (parent_doc_text, child_doc_list)


def process_image(image_path, metadata):
    """
    Exracts data from the image and returns the data from the image and list of child chunks 
    """
    image_data = preprocess_image_on_size_and_words(image_path)
    child_doc_list = []
    if not image_data:
        # Generate questions and answers via llm
        image_data, child_doc_list = process_image_via_llm(
            image_path, metadata)

    return (image_data, child_doc_list)


def index_doc_for_parent_child_retriever(docs):
    """
    Accepts a list of Document(s) and creates vectorstore and doc store for it
    """
    parent_doc_list = []
    child_doc_list = []
    ppt_title = None
    last_narrative_text = None
    larger_chunk = ""  # THIS CREATES CHILD CHUNKS, SUMMARY AND QUESTIONS
    doc_id = str(uuid.uuid4())
    # THIS GETS STORED IN THE DOC STORE, THE KEY DIFFERENCE IS THIS INCLUDES ALL THE CONTENT FROM A SLIDE
    parent_doc_text = ""
    for doc in docs:
        page_content = doc.page_content
        metadata = doc.metadata
        category = metadata['category']
        metadata['id'] = doc_id
        table = metadata.get('text_as_html')
        image_path = metadata.get('image_path')

        if category == 'Title':
            if larger_chunk:
                parent_doc_list.append((metadata['id'], Document(
                    page_content=parent_doc_text, metadata=metadata)))  # hold the metadata for it?
                # TODO: THE ID STORED IN METADATA IS IN CORRECT AS THE DICT IS BEING PASSED BY REFERENCE.
                # THIS DOES NOT EFFECT THE FUNCTIONALITY AS OF NOW. RESOLVE THIS
                # Generate the summary  and questions for larger chunk only
                if paragraph_parser.get_word_count(larger_chunk) > 50:
                    generate_summary = paragraph_parser.get_word_count(
                        larger_chunk) > 150
                    child_doc_list.extend(open_ai_util.ppt_slide_parser(
                        metadata=metadata,
                        ppt_title=ppt_title,
                        slide_text=larger_chunk,
                        generate_summary=generate_summary
                    ))
            doc_id = str(uuid.uuid4())  # update the doc id
            metadata['id'] = doc_id
            child_doc_list.append(Document(page_content=page_content, metadata={
                                  'doc_id': doc_id}))  # indexing the current title
            larger_chunk = f'{page_content}:\n'  # pick TITLE as well
            parent_doc_text = larger_chunk
            ppt_title = page_content
        elif page_content:
            larger_chunk += page_content
            # last_narrative_text = page_content
            parent_doc_text += page_content
            child_doc_list.extend(
                paragraph_parser.parse_paragraph(page_content, metadata))
        elif table:
            # metadata['id'] = str(uuid.uuid4())
            # WE SHOULD ALSO MAP THE TITLE AS CHILD AND RETURN THE TABLE FOR IT
            # THERE WILL BE TWO CHUNKS WITH TITLE ONE WILL POINT TO THE HTML TABLE AND ANOTHER TO SLIDE CONTENT
            # TODO EXPLORE HOW TO STORE THE ENTIRE SLIDE AS A SINGLE ELEMENT?
            # if ppt_title:
            #     child_doc_list.append(Document(page_content=ppt_title, metadata={
            #         'doc_id': metadata['id']}))
            # DECIDED TO STORE ENTIRE SLIDE AS ONE ENTITY
            # parent_doc_list.append((
            #     metadata['id'],
            #     Document(page_content=table, metadata=metadata)
            # ))
            parent_doc_text += f'\n{table}\n'
            child_doc_list.extend(open_ai_util.get_table_description(
                table=table, header=ppt_title, metadata=metadata))
        elif image_path:
            # # DO NOT PARSE IMAGES AS OF NOW
            image_data, image_child_chunks = process_image(
                image_path, metadata)
            # larger_chunk += image_data # larger_chunk #WE DO NOT WANT TO GENERATE SUMMARIES FOR THE IMAGES AS WELL AS OF NOW IN THE PARENT DOC AS WELL
            parent_doc_text += f'\n{image_data}\n'
            # Handle case when child_chunks = None
            child_doc_list.extend(image_child_chunks)

        if len(larger_chunk) > 4000:
            child_doc_list.extend(open_ai_util.ppt_slide_parser(
                metadata=metadata,
                ppt_title=ppt_title,
                slide_text=larger_chunk,
                generate_summary=True
            ))
            parent_doc_list.append((metadata['id'], Document(
                page_content=parent_doc_text, metadata=metadata)))
            larger_chunk = ""
            parent_doc_text = ""
            doc_id = str(uuid.uuid4())  # updated the uuid
    if larger_chunk:
        metadata = {'id': doc_id}
        parent_doc_list.append((metadata['id'], Document(
            page_content=parent_doc_text, metadata=metadata)))  # hold the metadata for it?
        # Generate the summary  and questions for larger chunk only
        if paragraph_parser.get_word_count(larger_chunk) > 50:
            generate_summary = paragraph_parser.get_word_count(
                larger_chunk) > 150
            child_doc_list.extend(open_ai_util.ppt_slide_parser(
                metadata=metadata,
                ppt_title=ppt_title,
                slide_text=larger_chunk,
                generate_summary=generate_summary
            ))
    with open('child_list_ppt.txt', 'w') as f:
        for d in child_doc_list:
            f.write(f'{d.page_content}\n\n')
    with open('parent_list_ppt.txt', 'w') as f:
        for _, d in parent_doc_list:
            f.write(f'{d.page_content}\n\n')
    ChromaDB().create_persistent_vector_database(docs=child_doc_list)
    store.mset(parent_doc_list)


def update_metadata(docs, file_path):
    # Further modifications could be done here.
    for doc in docs:
        doc.metadata['source'] = file_path
    return docs


def process(file_path):
    docs = process_pptx_files(file_path)
    docs = update_metadata(docs, file_path)
    index_doc_for_parent_child_retriever(docs)
