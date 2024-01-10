from PIL import Image
import sys
from langchain.storage import LocalFileStore
from langchain.storage._lc_store import create_kv_docstore
# from rag_fusion.chain import chain as rag_fusion_chain
import os
import openpyxl
import json
import functools
import shutil
from datetime import datetime, timedelta
import time
import uuid

from langchain.chains import RetrievalQA, HypotheticalDocumentEmbedder
from langchain.callbacks import get_openai_callback
from langchain.chat_models import ChatOpenAI
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain_community.vectorstores import SupabaseVectorStore
from supabase.client import Client, create_client
from langchain.retrievers import ParentDocumentRetriever
from langchain.storage import InMemoryStore
from langchain.text_splitter import *
from langchain.vectorstores.faiss import FAISS
from langchain.docstore.document import Document
from langchain.document_loaders import UnstructuredFileLoader, UnstructuredURLLoader
import tensorflow as tf
import numpy as nf
import tensorflow_hub as hub
from langchain_core.embeddings import Embeddings

import test_llm_basics.open_ai_util as open_ai_integration
import supabase_docstore as custom_supabase
import paragraph_parser as custom_text_parser
import custom_embeddings


from dotenv import load_dotenv
load_dotenv()


llm_model_type = "gpt-4-1106-preview"
llm = ChatOpenAI(max_retries=3, model=llm_model_type)
embeddings = custom_embeddings.CustomFaissEmbeddings()
openai_embeddings = OpenAIEmbeddings()
company_handbook_faiss_path = "./faiss_company_handbook"

# overiding defaults from env variables
documents_table = os.environ.get('SUPABASE_DOCUMENTS_TABLE', '')
embeddings_table = os.environ.get('SUPABASE_EMBEDDINGS_TABLE', '')
store = custom_supabase.SupabaseDocstore(documents_table)


def create_persistent_local_doc_store(store_location='./store_location'):
    """
    Creates a persistent doc store that does not erase even when the 
    """
    fs = LocalFileStore(store_location)
    store = create_kv_docstore(fs)
    return store


def create_local_doc_store():
    """
    Returns an instance of in memory doc store
    """
    return InMemoryStore()


def create_supabase_vectorstore(docs, table_name, use_local_embedding=False):
    """
    Creates a vectorstore on supabase with <docs> using the <table_name> table 
    """
    supabase_url = os.environ.get("SUPABASE_URL")
    supabase_key = os.environ.get("SUPABASE_KEY")
    supabase: Client = create_client(supabase_url, supabase_key)
    embeddings = custom_embeddings.CustomSupabaseEmbeddings(
    ) if use_local_embedding else openai_embeddings

    vectorstore = SupabaseVectorStore.from_documents(
        docs,
        embeddings,
        client=supabase,
        table_name=table_name
    )
    return vectorstore


def initialize_vectorstore(input, supabase):
    """
    This function initializes a vector store using FAISS from input text documents and embeddings.
    """
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=700, chunk_overlap=200)
    texts = ""
    vectorstore = None
    if (input is not None):
        texts = text_splitter.split_documents(input)
        for text in texts:
            print(text)
    if supabase:
        vectorstore = create_supabase_vectorstore(texts)
    else:
        vectorstore = FAISS.from_documents(texts, embeddings)
    print("Vectorstore created.")

    return vectorstore


def remove_folder(path):
    try:
        shutil.rmtree(path)
        print(f"Successfully removed folder at {path}")
    except Exception as e:
        print(f"Error while removing folder at {path}: {e}")


def save_faiss_locally(vectorstore, path: str):
    """
    This function saves a vectorstore locally at a specified path and prints a confirmation message.

    """
    vectorstore.save_local(path)  # type: ignore

    print(f"Vectorstore saved locally")


def get_vectorstore_retriever(index_path: str):
    """
    This function returns a retriever object for a vector store index, loaded from a path (string) using
    Meta's FAISS library.
    """
    return FAISS.load_local(index_path, embeddings=embeddings).as_retriever()


def merge_with_old_vectorstore(new_vectorstore):
    """
    This function merges a new vectorstore with an old vectorstore, saves the merged vectorstore
    locally, and removes the old vectorstore.

    """
    old_vectorstore = None
    try:
        old_vectorstore = FAISS.load_local(
            company_handbook_faiss_path, embeddings)
    except Exception as e:
        print(e)
    if old_vectorstore:
        old_vectorstore.merge_from(new_vectorstore)
        print("Vectorstores merged.")
        remove_folder(company_handbook_faiss_path)

    save_faiss_locally(old_vectorstore or new_vectorstore,
                       company_handbook_faiss_path)

    return old_vectorstore or new_vectorstore


def format_documents(quick_connectors_data):
    """
    ## CURRENTLY NOT IN USE, SHALL BE UPDATED LATER  
    """
    document = Document(page_content=quick_connectors_data.pop(
        'text'), metadata=quick_connectors_data)
    return document


def process_excel_file(file_path):
    """
    Reads excel manually and documents them in form of sentences
    ## CURRENTLY NOT IN USE, SHALL BE UPDATED LATER
    """
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    sheet = workbook.active
    tmp_file = 'tmp_transformed.txt'
    file = open(tmp_file, 'w')
    for row in sheet.iter_rows(min_row=2, max_col=sheet.max_column, max_row=sheet.max_row):
        row_data = [cell.value for cell in row]
        name, marks, sub = row_data[1:]
        file.write(
            f'Student with name {name} got {marks} marks in the subject {sub}')

    return tmp_file


def process_csv_files(csv_file_path):
    """
    ## CURRENTLY NOT IN USE, SHALL BE UPDATED LATER
    """
    import csv
    is_header = True
    file = open(f'tmp_transformed.txt', 'w')
    with open(csv_file_path, 'r') as csvfile:
        csvreader = csv.reader(csvfile)
        for row_data in csvreader:
            area, date, dealer, sales = row_data
            if not is_header:
                file.write(
                    f'On {date}, sales in the {area} area by the dealer {dealer} were {sales}. \n')
            is_header = False
    return f'tmp_transformed.txt'


def create_and_merge(path, mode='single'):
    """
    creates a vector store and merges it with existing vector store
    # CURRENTLY NOT IN USE, SHALL BE UPDATED LATER
    """
    # _, file_extension = os.path.splitext(path)
    # if file_extension.lower() == ".csv":
    #     path = process_csv_files(path)
    # elif file_extension.lower() == ".xlsx":
    #     path = process_excel_file(path)
    # print(file_extension)
    # print(path)
    document = UnstructuredFileLoader(
        file_path=path,
        strategy='hi_res',
        mode=mode,
        skip_infer_table_types=[],
        pdf_infer_table_structure=True,
        pdf_extract_images=True
    ).load()
    new_vectorstore = initialize_vectorstore(document, False)
    return merge_with_old_vectorstore(new_vectorstore)


def process_pptx_files(filepath_list):
    """
    Loads list pptx files in the form of key value pairs 
    ## CURRENTLY NOT IN USE, SHALL BE UPDATED LATER
    """
    from pptx import Presentation
    file_dict = {}
    for file in filepath_list:
        f = open(f'{file}', "rb")
        prs = Presentation(f)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_image:
                    with open('test.jpeg', 'w') as f:
                        f.write(shape.image)
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        file_dict[file] = " ".join(text_runs)
    return file_dict


def process_extracted_files(image_path_list):
    """
    Process extracted image files and created parent and child chunks on that
    """
    parent_doc_list = []
    child_doc_list = []

    for metadata in image_path_list:
        last_title = metadata['last_title']
        image_path = metadata['image_path']
        last_narrative_text = metadata['last_narrative_text']
        page_content = open_ai_integration.get_image_description(
            image_path, last_title, last_narrative_text)
        # CHECK FOR A VALID RESPONSE
        if page_content not in ['UNEXTRACTABLE_DATA', 'UNPROCESSABLE_ENTITY']:
            image_id = str(uuid.uuid4())
            # use the same metadata as of the parent image
            metadata['id'] = image_id
            parent_doc_list.append(
                (image_id, Document(page_content=page_content, metadata=metadata)))
            child_doc_list.extend(
                custom_text_parser.parse_paragraph(page_content, metadata))
            child_doc_list.extend(
                open_ai_integration.get_paragraph_description(page_content, metadata))
    return (parent_doc_list, child_doc_list)


def rectangle_area(coordinates):
    """"
    Returns area of a rectangle
    """
    (x1, y1), (x2, y2), (x3, y3), (x4, y4) = coordinates
    width = max(x1, x2, x3, x4) - min(x1, x2, x3, x4)
    height = max(y1, y2, y3, y4) - min(y1, y2, y3, y4)
    area = width * height
    return area


def is_intersecting(rect1, rect2):
    """
    Checks whether or not two rectangles intersect
    """
    l1, _, r1, _ = rect1
    l2, _, r2, _ = rect2
    x1, y1 = l1
    x2, y2 = r1
    x3, y3 = l2
    x4, y4 = r2
    if x1 > x4 or x3 > x2:
        return False
    if y1 > y4 or y3 > y2:
        return False
    return True


def is_overlapping(parent, child):
    """
    Checks whether or not two rectangles overlap
    """
    l1, _, r1, _ = parent
    l2, _, r2, _ = child
    x1, y1 = l1
    x2, y2 = r1
    x3, y3 = l2
    x4, y4 = r2
    if x1 <= x3 and y1 <= y3 and x2 >= x4 and y2 >= y4:
        return True
    return False


def compare(item1, item2):
    """
    Sorts cordinates on the basis of start point and area
    """
    cord_1, _, _, _ = item1
    cord_2, _, _, _ = item2
    x1, y1 = cord_1
    x2, y2 = cord_2
    if x1 == x2 and y1 == y2:
        a1 = rectangle_area(item1)
        a2 = rectangle_area(item2)
        return a2-a1
    elif (x1 <= x2 and y1 <= y2):
        return -1
    else:
        return 1


def process_image_overlaps(cor_list, image_path_mapping):
    """
    Eliminates all the child images that unstructured generates during parsing and retains only the parent image.
    """

    cor_list.sort(key=functools.cmp_to_key(compare))
    final_images_cords = {}

    for x in cor_list:
        def f(x):
            (x1, y1), _, (x2, y2), _ = x
            return f'{x1}_{y1}_{x2}_{y2}'
        (x1, y1), _, _, _ = x
        key = f'{x1}_{y1}'
        new_rect = True
        curr_page = image_path_mapping[f(x)]['page_number']
        for existing_rectangle in final_images_cords.values():
            tmp_page = image_path_mapping[f(existing_rectangle)]['page_number']
            # if curr_page == tmp_page and is_intersecting(existing_rectangle, x):
            if curr_page == tmp_page and is_overlapping(existing_rectangle, x):
                new_rect = False
                break
        if new_rect:
            final_images_cords[key] = x
    final_image_path_list = []
    # the number of final images in a doc will be handful so this loop should not be much work
    # shall look in future on improvising this
    for final in final_images_cords.values():
        (x1, y1), _, (x2, y2), _ = final
        final_image_path_list.append(
            image_path_mapping[f'{x1}_{y1}_{x2}_{y2}'])

    return final_image_path_list


def create_parent_child_vectorstore(file_path, use_local_vectorstore=False, use_local_embeddings=False):
    """
    Takes a file path, creates child chunks and parent documents and loads them to vectorstore and doc store respectively
    ## WORKS WELL ONLY FOR PDF AS OF NOW
    """
    doc = UnstructuredFileLoader(
        file_path=file_path,
        strategy='hi_res',
        mode='elements',
        skip_infer_table_types=[],
        pdf_infer_table_structure=True,
        pdf_extract_images=True,
        chunking_strategy='by-title'
    ).load()
    cor_list = []
    image_path_mapping = {}
    SUB_TEXT_LIST = ['Text', 'NarrativeText',
                     'BulletedText']  # Improvise this later
    child_document_list = []
    # TODO: Update this to parent_document_list.
    larger_document_list = []
    final_image_path_list = []
    last_title = None
    last_narrative_text = None
    curr_chunk = ""
    doc_id = str(uuid.uuid4())
    for d in doc:
        page_content = d.page_content
        category = d.metadata['category']
        d.metadata['id'] = doc_id
        table_details = d.metadata.get('text_as_html')
        # NARRATIVE TEXT HAS TO be one of text, narrative_text, BulletedText
        if category in SUB_TEXT_LIST:
            last_narrative_text = page_content
        # SKIP INFERING PAGE NUMBERS
        if page_content.isdigit():
            continue
        # PROCESSING BEGINS

        # PROCESSING TITLE
        if d.metadata.get('category') == 'Title':
            if last_title and len(curr_chunk) > 0:
                larger_document_list.append((
                    d.metadata['id'],
                    Document(
                        page_content=f'{last_title}\n{curr_chunk}', metadata=d.metadata)
                ))
            doc_id = str(uuid.uuid4())
            curr_chunk = ""
            last_title = f'{d.page_content}: \n'
        # PROCESSING IMAGE
        elif category == 'Image':
            image_path = d.metadata.get('image_path')
            cor_list.append(d.metadata['coordinates']['points'])
            (x1, y1), _, (x2, y2), _ = d.metadata['coordinates']['points']
            d.metadata['last_title'] = last_title
            d.metadata['last_narrative_text'] = last_narrative_text
            image_path_mapping[f'{x1}_{y1}_{x2}_{y2}'] = d.metadata
            # DELAY PROCESSING OF IMAGE
        # PROCESSING TABLE
        elif table_details:
            # TODO: Evaluate if providing title and narrative text here improves accuracy.
            table_uuid = str(uuid.uuid4())
            d.metadata['id'] = table_uuid
            larger_document_list.append((
                table_uuid,
                Document(
                    page_content=table_details, metadata=d.metadata)
            ))
            # PASSING IN LAST_TITLE AND NARRATIVE TEXT FOR TABLES AS WELL
            child_document_list.extend(
                open_ai_integration.get_table_description(table_details, d.metadata, last_title, last_narrative_text))
        else:
            child_document_list.extend(custom_text_parser.parse_paragraph(
                d.page_content, d.metadata, last_title))

            if len(curr_chunk) + len(page_content) > 4000:
                # update meta deta, currently its inappropriate
                # HANDLE paragraphs pre-chunked to multiple pages
                larger_document_list.append((d.metadata['id'],
                                             Document(page_content=f'{last_title}\n{curr_chunk}', metadata=d.metadata)))
                doc_id = str(uuid.uuid4())  # update the doc_id for new parent
                curr_chunk = page_content
            else:
                curr_chunk += (page_content)

            # list of documents with summarised paragraph and questions if the size is greater than 300
            add_summary = len(page_content) > 300
            if add_summary:
                child_document_list.extend(
                    open_ai_integration.get_paragraph_description(page_content, d.metadata, add_summary))

    if len(curr_chunk) > 0:
        larger_document_list.append((doc_id,
                                     Document(page_content=f'{last_title}\n{curr_chunk}', metadata={'id': doc_id})))
        if len(curr_chunk) > 300:  # temporarily do not summarise small paragraphs
            child_document_list.extend(
                open_ai_integration.get_paragraph_description(curr_chunk, d.metadata, len(curr_chunk) > 300))

        # child_document_list.extend(custom_text_parser.parse_paragraph(
        #     d.page_content, d.metadata, last_title))
        # d.page_content = f'{last_title}{d.page_content}'
    final_image_path_list = process_image_overlaps(
        cor_list, image_path_mapping)
    print(final_image_path_list)
    if final_image_path_list:
        a, b = process_extracted_files(final_image_path_list)
        print("RESPONSE FROM IMAGES--> parent ", a)
        print("RESPONSE FROM IMAGES--> child ", b)
        child_document_list.extend(b)
        larger_document_list.extend(a)
    # # TODO: improve this to pass in parent document

    # Writes down all the chunks into child_list.txt file
    # Helpful for debugging
    with open('child_list.txt', 'w') as f:
        for d in child_document_list:
            f.write(f'{d.page_content}\n\n')

    vectorstore = None
    if use_local_vectorstore:
        vectorstore = FAISS.from_documents(
            documents=child_document_list, embedding=embeddings)
        merge_with_old_vectorstore(vectorstore)
    else:
        vectorstore = create_supabase_vectorstore(
            child_document_list, embeddings_table, use_local_embeddings)  # using the embeddings table picked from env var

    store.mset(larger_document_list)
    return vectorstore


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python loader.py file_path1 file_path2 ...")
        exit(1)
    file_paths = sys.argv[1:]
    for file_path in file_paths:
        create_parent_child_vectorstore(file_path)
