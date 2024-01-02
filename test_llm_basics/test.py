import openpyxl
import json
import shutil
from datetime import datetime, timedelta
import time

from langchain.chains import RetrievalQA
from langchain.callbacks import get_openai_callback
from langchain.chat_models import ChatOpenAI
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.text_splitter import *
from langchain.vectorstores.faiss import FAISS
from langchain.docstore.document import Document
from langchain.document_loaders import UnstructuredFileLoader
import tensorflow as tf
import tensorflow_hub as hub
from langchain_core.embeddings import Embeddings


class CustomEmbeddings(Embeddings):
    """
    class to utilise UniversalSentenceEncoder for embeddings
    """
    MODEL_URL = "./models/UniversalSentenceEncoder"

    def load_model(self):
        model = hub.load(self.MODEL_URL)
        return model

    def embed_query(self, input):
        model = self.load_model()
        return model([input])[0]

    def embed_documents(self, input):
        model = self.load_model()
        return model(input)


llm_model_type = "gpt-3.5-turbo"
openai_api_key = "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"
llm = ChatOpenAI(max_retries=3, openai_api_key=openai_api_key, temperature=0, model=llm_model_type)
# embeddings = OpenAIEmbeddings(openai_api_key=openai_api_key, max_retries=1)
embeddings = CustomEmbeddings()
company_handbook_faiss_path = "./faiss_company_handbook"


def initialize_vectorstore(input):
    """
    This function initializes a vector store using FAISS from input text documents and embeddings.
    """
    text_splitter = RecursiveCharacterTextSplitter()
    texts = ""
    if (input is not None):
        texts = text_splitter.split_documents(input)
        for text in texts:
            print(text)
    succeeded = False
    counter = 1
    while not succeeded:
        try:
            vectorstore = FAISS.from_documents(texts, embeddings)  # type: ignore
            succeeded = True
        except Exception as e:
            print(e)
            time.sleep(65)
            counter += 1

    print("Vectorstore created.")

    return vectorstore


def remove_folder(path):
    try:
        shutil.rmtree(path)
        print(f"Successfully removed folder at {path}")
    except Exception as e:
        print(f"Error while removing folder at {path}: {e}")


def save_faiss_locally(vectorstore, path: str):
    """
    This function saves a vectorstore locally at a specified path and prints a confirmation message.

    """
    vectorstore.save_local(path)  # type: ignore

    print(f"Vectorstore saved locally")


def get_vectorstore_retriever(index_path: str):
    """
    This function returns a retriever object for a vector store index, loaded from a path (string) using
    Meta's FAISS library.
    """
    return FAISS.load_local(index_path, embeddings=embeddings).as_retriever()


def merge_with_old_vectorstore(new_vectorstore):
    """
    This function merges a new vectorstore with an old vectorstore, saves the merged vectorstore
    locally, and removes the old vectorstore.

    """
    old_vectorstore = None
    try:
        old_vectorstore = FAISS.load_local(
            company_handbook_faiss_path, embeddings)
    except Exception as e:
        print(e)
    if old_vectorstore:
        old_vectorstore.merge_from(new_vectorstore)
        print("Vectorstores merged.")
        remove_folder(company_handbook_faiss_path)

    save_faiss_locally(old_vectorstore or new_vectorstore, company_handbook_faiss_path)

    return old_vectorstore or new_vectorstore


def format_documents(quick_connectors_data):
    document = Document(page_content=quick_connectors_data.pop('text'), metadata=quick_connectors_data)
    return document


def upload_files():
    quick_connectors_docx_data = """
    {
      "editedBy": "Ekansh Verma",
      "id": "16MX6uXUmhwXKNqTd2enzTs-thMV_9bs5m5RCiPeUiU0",
      "mimeType": "application/vnd.google-apps.document",
      "modifiedTime": "2023-12-27T10:07:17.997Z",
      "text": "Ekansh Verma:\\r\\n* Age: 17\\r\\n* Height: 5 feet 9 inches\\r\\n* Hobbies: Avid reader, guitar player, passion for exploring different genres of music.\\r\\n* Co-curricular Activities: Member of the school debate team, actively involved in the photography club.\\r\\nVikram:\\r\\n* Age: 16\\r\\n* Height: 6 feet 1 inch\\r\\n* Hobbies: Sports enthusiast with a keen interest in basketball, enjoys sketching and pencil/charcoal drawings.\\r\\n* Co-curricular Activities: Member of the school's science club, serves as the captain of the basketball team.\\r\\nVarun:\\r\\n* Age: 18\\r\\n* Height: 5 feet 10 inches\\r\\n* Hobbies: Technology enthusiast, coding and building applications, passion for robotics, enjoys strategic board games.\\r\\n* Co-curricular Activities: Active member of the school's robotics club, contributes to the school's coding competition team.\\r\\nEshank Verma:\\r\\n* Age: 17\\r\\n* Height: 5 feet 7 inches\\r\\n* Hobbies: Flair for the arts, particularly theater, enjoys acting, passion for gardening.\\r\\n* Co-curricular Activities: Key member of the drama club, actively involved in the school's environmental club.",
      "title": "HEllo",
      "url": "https://docs.google.com/document/d/16MX6uXUmhwXKNqTd2enzTs-thMV_9bs5m5RCiPeUiU0/edit?usp=drivesdk"
    }
    """
    quick_connectors_sheet_data = """
    {
      "editedBy": "Ekansh Verma",
      "id": "1TmfbWggucsKB27LtqRsgLIm6p9khQxdFw9dAkDtb4HY",
      "mimeType": "application/vnd.google-apps.spreadsheet",
      "modifiedTime": "2023-12-27T09:46:51.893Z",
      "text": "Id,name,Marks,Subject\\n1,Ekansh Verma,80,Maths\\n2,Ekansh Verma,65,English\\n3,Ekansh Verma,73,Phy\\n4,Ekansh Verma,59,Chem\\n5,Vikram,80,Maths\\n6,Vikram,75,English\\n7,Vikram,89,Phy\\n8,Vikram,99,Chem\\n9,Varun,14,Maths\\n10,Varun,36,English\\n11,Varun,54,Phy\\n12,Varun,65,Chem\\n13,Eshank Verma,59,Maths\\n14,Eshank Verma,66,English\\n15,Eshank Verma,72,Phy\\n16,Eshank Verma,80,Chem",
      "title": "Untitled spreadsheet",
      "url": "https://docs.google.com/spreadsheets/d/1TmfbWggucsKB27LtqRsgLIm6p9khQxdFw9dAkDtb4HY/edit?usp=drivesdk"
    }
    """

    def create_and_merge(path):
        """
        creates a vector store and merges it with existing vector store
        """
        document = UnstructuredFileLoader(path).load()
        new_vectorstore = initialize_vectorstore(document)
        return merge_with_old_vectorstore(new_vectorstore)
    # document = format_documents(json.loads(quick_connectors_docx_data))
    # document = UnstructuredFileLoader('Student Marks.xlsx').load()
    # new_vectorstore = create_and_merge('presentation (11.12.2022).pptx')
    new_vectorstore = None
    for i in range(1, 6):
        print("loading")
        new_vectorstore = create_and_merge(f'img{i}.png')
    # create_and_merge('Box Cricket Rules (JPL 2023).docx')
    # create_and_merge('Career Ladder Specs and Feedback forms - Oct 2023.docx')
    # new_vectorstore = create_and_merge('Django Setup Instructions.docx')
    # vectorstore = merge_with_old_vectorstore(initialize_vectorstore([document]))
    # document = UnstructuredFileLoader(
    # 'Student Marks.xlsx').load()
    # new_vectorstore.max_marginal_relevance_search
    # document = format_documents(json.loads(quick_connectors_sheet_data))
    # new_vectorstore.similarity_search()
    ask_questions(new_vectorstore)


def ask_questions(vector_store):
    """
    Selects answers from the vectors store and writes them to ouput.txt file
    """
    # ques_list = ["Hobbies of EKansh", "Total Students in the list",
    #                  "Hobbies comparison of Ekansh and Vikram", "Name of the oldest student in the list"]
    # sheet_ques_list = ["Maths", "Ekansh", "marks in Maths", "marks of Ekansh", "marks of name Ekansh", "marks of Vikram", "marks in phy", "marks in PHY",
    #                    "sum of marks of Ekansh"]
    jpl_questions = ["who are the members of the project?"]
    with open('output.txt', 'w') as f:
        for ques in jpl_questions:
            answers = vector_store.similarity_search(query=ques, k=4)
            f.write(f"Answer for {ques} is \n")
            d = 1
            for answer in answers:
                f.write(f"Response number: {d}")
                d += 1
                f.write(f"\n {answer.page_content} \n")
            print("answered ", ques)
    # qa = RetrievalQA.from_llm(
    #     llm=llm, retriever=vector_store.as_retriever(), return_source_documents=True)

    # with get_openai_callback() as cb:
    # result = qa({"query": "Ekansh Verma"})

    # Preliminary formatted results from chain
    # chat_answer = result["result"]


upload_files() # loads vector store and answers questions


def read_excel_file(file_path):
    """
    Reads excel manually and documents them in form of sentences
    """
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    sheet = workbook.active
    documents = []
    data = {}
    for row in sheet.iter_rows(min_row=2, max_col=sheet.max_column, max_row=sheet.max_row):
        row_data = [cell.value for cell in row]
        name, marks, sub = row_data[1:]
        default_dict = data.get(name, {})
        default_dict['subjects'] = default_dict.get('subjects', {})
        default_dict['subjects'][sub] = marks
        data[name] = default_dict

    for key, value in data.items():
        documents.append(Document(
            page_content=json.dumps({key: value}),
            metadata={
                "Student Name": row_data[1], "Subject": row_data[3], "marks": row_data[2]}
        ))

    print(documents)
    return documents

# Replace 'your_file.xlsx' with the actual path to your Excel file
# file_path = 'Student Marks.xlsx'
# vectorstore = FAISS.from_documents(read_excel_file(file_path), embeddings)
# ask_questions(vectorstore)


def process_pptx_files(filepath_list):
    """
    Loads list pptx files in the form of key value pairs
    """
    from pptx import Presentation
    file_dict = {}
    for file in filepath_list:
        f = open(f'{file}', "rb")
        prs = Presentation(f)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        file_dict[file] = " ".join(text_runs)
    return file_dict
